"""
Phase 1 — Extractor Tests (PDF, PPTX, plain-text, toggle)
===========================================================
Tests the new rich-format ingestion layer introduced in:
  core/ingestion/extractors.py
  core/ingestion/loader.load_document()
  core/ingestion/ingestor.Ingestor.ingest_file()

Tests are self-contained: they create and clean up all fixture files in a
temporary directory — no network access, no running services required.

Run:
    pytest tests/phase1/test_extractors.py -v
"""
from __future__ import annotations

import tempfile
from pathlib import Path

import pytest

from core.ingestion.extractors import (
    PdfExtractor,
    PptxExtractor,
    TextExtractor,
    extract_text,
    get_extractor,
)
from core.ingestion.loader import load_document, load_text_file


# ── Helpers / shared fixtures ─────────────────────────────────────────────────

@pytest.fixture(scope="module")
def tmp_dir():
    """Module-scoped temporary directory; cleaned up after all tests in module."""
    with tempfile.TemporaryDirectory() as d:
        yield Path(d)


@pytest.fixture(scope="module")
def txt_file(tmp_dir: Path) -> Path:
    """Plain-text fixture file."""
    p = tmp_dir / "sample.txt"
    p.write_text(
        "To be or not to be, that is the question.\n"
        "Whether 'tis nobler in the mind to suffer.",
        encoding="utf-8",
    )
    return p


@pytest.fixture(scope="module")
def pdf_file(tmp_dir: Path) -> Path:
    """
    Minimal PDF fixture with extractable text, created via pypdf.
    Requires: pypdf>=4.0 (already in requirements.txt).
    """
    from pypdf import PdfWriter
    from pypdf.generic import NameObject, DictionaryObject, ArrayObject, DecodedStreamObject

    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)

    # Inject a minimal content stream with readable text
    content_stream = DecodedStreamObject()
    content_stream.set_data(b"BT /F1 12 Tf 100 700 Td (PolyRAG PDF ingestion test.) Tj ET")
    page[NameObject("/Contents")] = writer._add_object(content_stream)  # type: ignore[attr-defined]

    resources = DictionaryObject()
    font_dict = DictionaryObject()
    font_entry = DictionaryObject()
    font_entry[NameObject("/Type")] = NameObject("/Font")
    font_entry[NameObject("/Subtype")] = NameObject("/Type1")
    font_entry[NameObject("/BaseFont")] = NameObject("/Helvetica")
    font_dict[NameObject("/F1")] = writer._add_object(font_entry)  # type: ignore[attr-defined]
    resources[NameObject("/Font")] = font_dict
    page[NameObject("/Resources")] = resources

    p = tmp_dir / "sample.pdf"
    with p.open("wb") as f:
        writer.write(f)
    return p


@pytest.fixture(scope="module")
def pptx_file(tmp_dir: Path) -> Path:
    """
    Real PPTX fixture with title, body text, and speaker notes.
    Requires: python-pptx>=1.0 (already in requirements.txt).
    """
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title + Content layout
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = "PolyRAG PPTX Test Slide"
    slide.placeholders[1].text = "PDF and PowerPoint ingestion is fully supported."

    notes = slide.notes_slide.notes_text_frame
    notes.text = "Speaker note: rich-format ingestion test."

    p = tmp_dir / "sample.pptx"
    prs.save(str(p))
    return p


@pytest.fixture(scope="module")
def legacy_ppt_file(tmp_dir: Path) -> Path:
    """Dummy .ppt file path (we never write real content — just test extension guard)."""
    p = tmp_dir / "legacy.ppt"
    p.write_bytes(b"\xd0\xcf\x11\xe0")  # OLE2 magic bytes — not a real ppt, just for path
    return p


# ── TextExtractor ─────────────────────────────────────────────────────────────

class TestTextExtractor:
    def test_extracts_content(self, txt_file: Path):
        ext = TextExtractor()
        result = ext.extract(txt_file)
        assert "To be or not to be" in result

    def test_preserves_full_content(self, txt_file: Path):
        ext = TextExtractor()
        result = ext.extract(txt_file)
        original = txt_file.read_text(encoding="utf-8")
        assert result == original

    def test_known_extensions(self):
        assert ".txt" in TextExtractor.extensions
        assert ".md" in TextExtractor.extensions
        assert ".csv" in TextExtractor.extensions


# ── PdfExtractor ──────────────────────────────────────────────────────────────

class TestPdfExtractor:
    def test_raises_on_missing_pypdf(self, monkeypatch):
        """If pypdf is not installed, ImportError with instructions is raised."""
        import builtins
        real_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "pypdf":
                raise ImportError("No module named 'pypdf'")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", mock_import)
        ext = PdfExtractor()
        with pytest.raises(ImportError, match="pip install"):
            ext.extract(Path("dummy.pdf"))

    def test_raises_on_blank_pdf(self, tmp_dir: Path):
        """A PDF with no text content raises ValueError with a clear message."""
        from pypdf import PdfWriter
        blank = tmp_dir / "blank.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        with blank.open("wb") as f:
            writer.write(f)

        ext = PdfExtractor()
        with pytest.raises(ValueError, match="No extractable text"):
            ext.extract(blank)

    def test_extension(self):
        assert ".pdf" in PdfExtractor.extensions


# ── PptxExtractor ─────────────────────────────────────────────────────────────

class TestPptxExtractor:
    def test_extracts_title(self, pptx_file: Path):
        ext = PptxExtractor()
        result = ext.extract(pptx_file)
        assert "PolyRAG PPTX Test Slide" in result

    def test_extracts_body(self, pptx_file: Path):
        ext = PptxExtractor()
        result = ext.extract(pptx_file)
        assert "PowerPoint ingestion" in result

    def test_extracts_notes(self, pptx_file: Path):
        ext = PptxExtractor()
        result = ext.extract(pptx_file)
        assert "Speaker note" in result

    def test_slide_title_prefix(self, pptx_file: Path):
        ext = PptxExtractor()
        result = ext.extract(pptx_file)
        assert "[Slide 1 Title]" in result

    def test_notes_prefix(self, pptx_file: Path):
        ext = PptxExtractor()
        result = ext.extract(pptx_file)
        assert "[Notes]" in result

    def test_raises_on_legacy_ppt(self, legacy_ppt_file: Path):
        ext = PptxExtractor()
        with pytest.raises(ValueError, match=r"\.ppt.*not supported"):
            ext.extract(legacy_ppt_file)

    def test_raises_on_missing_pptx(self, monkeypatch):
        """If python-pptx is not installed, ImportError with instructions is raised."""
        import builtins
        real_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "pptx":
                raise ImportError("No module named 'pptx'")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", mock_import)
        ext = PptxExtractor()
        with pytest.raises(ImportError, match="pip install"):
            ext.extract(Path("dummy.pptx"))

    def test_extensions(self):
        assert ".pptx" in PptxExtractor.extensions
        assert ".ppt" in PptxExtractor.extensions


# ── Extractor registry / get_extractor ───────────────────────────────────────

class TestExtractorRegistry:
    def test_routes_txt_to_text_extractor(self):
        assert isinstance(get_extractor(Path("doc.txt")), TextExtractor)

    def test_routes_md_to_text_extractor(self):
        assert isinstance(get_extractor(Path("README.md")), TextExtractor)

    def test_routes_pdf_to_pdf_extractor(self):
        assert isinstance(get_extractor(Path("report.pdf")), PdfExtractor)

    def test_routes_pptx_to_pptx_extractor(self):
        assert isinstance(get_extractor(Path("deck.pptx")), PptxExtractor)

    def test_routes_ppt_to_pptx_extractor(self):
        assert isinstance(get_extractor(Path("deck.ppt")), PptxExtractor)

    def test_unknown_extension_falls_back_to_text(self):
        assert isinstance(get_extractor(Path("notes.xyz")), TextExtractor)

    def test_case_insensitive_routing(self):
        assert isinstance(get_extractor(Path("REPORT.PDF")), PdfExtractor)
        assert isinstance(get_extractor(Path("DECK.PPTX")), PptxExtractor)


# ── extract_text (dispatcher with toggle) ────────────────────────────────────

class TestExtractText:
    def test_txt_enabled(self, txt_file: Path):
        result = extract_text(txt_file, enable_rich_formats=True)
        assert "To be or not to be" in result

    def test_txt_disabled(self, txt_file: Path):
        """Plain-text files are always accepted even when toggle is off."""
        result = extract_text(txt_file, enable_rich_formats=False)
        assert "To be or not to be" in result

    def test_pptx_enabled(self, pptx_file: Path):
        result = extract_text(pptx_file, enable_rich_formats=True)
        assert "PolyRAG PPTX Test Slide" in result

    def test_pptx_disabled_raises(self, pptx_file: Path):
        with pytest.raises(ValueError, match="enable_rich_formats=False"):
            extract_text(pptx_file, enable_rich_formats=False)

    def test_pdf_disabled_raises(self, tmp_dir: Path):
        dummy = tmp_dir / "dummy.pdf"
        dummy.write_bytes(b"%PDF-1.4")
        with pytest.raises(ValueError, match="enable_rich_formats=False"):
            extract_text(dummy, enable_rich_formats=False)

    def test_missing_file_raises_file_not_found(self, tmp_dir: Path):
        with pytest.raises(FileNotFoundError):
            extract_text(tmp_dir / "does_not_exist.txt")


# ── load_document (loader.py dispatcher) ─────────────────────────────────────

class TestLoadDocument:
    def test_txt_unchanged(self, txt_file: Path):
        """load_document for .txt returns same content as load_text_file."""
        via_new = load_document(txt_file)
        via_old = load_text_file(txt_file)
        assert via_new == via_old

    def test_pptx_returns_string(self, pptx_file: Path):
        result = load_document(pptx_file)
        assert isinstance(result, str)
        assert len(result) > 0

    def test_accepts_string_path(self, txt_file: Path):
        result = load_document(str(txt_file))
        assert "To be or not to be" in result

    def test_toggle_off_blocks_pptx(self, pptx_file: Path):
        with pytest.raises(ValueError, match="enable_rich_formats=False"):
            load_document(pptx_file, enable_rich_formats=False)


# ── End-to-end: ingest PPTX through the full pipeline ────────────────────────

class TestIngestorRichFormats:
    """Full ingest_file() integration test using in-memory ChromaDB."""

    @pytest.fixture(scope="class")
    def embedder(self):
        from core.embedding.sentence_transformer import SentenceTransformerProvider
        provider = SentenceTransformerProvider(
            {"model": "all-MiniLM-L6-v2", "device": "cpu", "batch_size": 32}
        )
        _ = provider.embedding_dim
        return provider

    @pytest.fixture(scope="class")
    def store(self):
        from core.store.registry import AdapterRegistry
        s = AdapterRegistry.create("chromadb", {"mode": "memory"})
        s.connect()
        yield s
        s.close()

    def _make_ingestor(self, store, embedder):
        from core.ingestion.ingestor import Ingestor
        cfg = {
            "collection_name": "extractor_test",
            "chunk_size": 300,
            "chunk_overlap": 30,
            "embed_batch_size": 8,
        }
        ing = Ingestor(store=store, embedder=embedder, config=cfg)
        return ing

    def test_ingest_txt_file(self, store, embedder, txt_file: Path):
        ing = self._make_ingestor(store, embedder)
        col = "test_ingest_txt"
        store.create_collection(col, embedder.embedding_dim)
        result = ing.ingest_file(str(txt_file), collection=col)
        assert result.upserted >= 1
        assert store.count(col) >= 1
        store.drop_collection(col)

    def test_ingest_pptx_file(self, store, embedder, pptx_file: Path):
        ing = self._make_ingestor(store, embedder)
        col = "test_ingest_pptx"
        store.create_collection(col, embedder.embedding_dim)
        result = ing.ingest_file(str(pptx_file), collection=col)
        assert result.upserted >= 1
        assert store.count(col) >= 1
        store.drop_collection(col)

    def test_ingest_pptx_metadata_contains_source_file(self, store, embedder, pptx_file: Path):
        """source_file metadata must be present on every chunk."""
        ing = self._make_ingestor(store, embedder)
        col = "test_meta_pptx"
        store.create_collection(col, embedder.embedding_dim)
        ing.ingest_file(str(pptx_file), collection=col)

        vec = embedder.embed_one("PolyRAG PPTX")
        results = store.query(col, vec, top_k=1)
        assert "source_file" in results[0].document.metadata
        store.drop_collection(col)

    def test_toggle_off_prevents_pptx_ingest(self, store, embedder, pptx_file: Path):
        ing = self._make_ingestor(store, embedder)
        with pytest.raises(ValueError, match="enable_rich_formats=False"):
            ing.ingest_file(str(pptx_file), enable_rich_formats=False)
