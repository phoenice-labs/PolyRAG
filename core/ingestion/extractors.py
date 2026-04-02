"""
Document extractors — convert PDF and PowerPoint files to plain text.

Each extractor returns a single ``str`` that is then handed to the existing
chunking / embedding / vector-store pipeline unchanged.

Toggle
------
Set ``ingestion.enable_rich_formats: false`` in config.yaml (or the env var
``POLYRAG_RICH_FORMATS=0``) to disable PDF/PPTX support and revert to
text-only ingestion without restarting the server.

Missing libraries
-----------------
If ``pypdf`` or ``python-pptx`` are not installed the extractor raises an
``ImportError`` with clear installation instructions.  Plain ``.txt`` files
are never affected.
"""
from __future__ import annotations

import os
from pathlib import Path


# ── Base ──────────────────────────────────────────────────────────────────────

class ExtractorBase:
    """Minimal interface every extractor must satisfy."""

    #: File extensions this extractor handles (lower-case, with dot).
    extensions: tuple[str, ...] = ()

    def extract(self, path: Path) -> str:
        """Return the full plain-text content of *path*."""
        raise NotImplementedError


# ── Plain-text (pass-through) ─────────────────────────────────────────────────

class TextExtractor(ExtractorBase):
    """Pass-through extractor for plain-text files."""

    extensions = (".txt", ".md", ".rst", ".csv", ".json", ".xml", ".html", ".htm")

    def extract(self, path: Path) -> str:
        return path.read_text(encoding="utf-8", errors="replace")


# ── PDF ───────────────────────────────────────────────────────────────────────

class PdfExtractor(ExtractorBase):
    """
    Extract text from PDF files using ``pypdf``.

    Install:  pip install "pypdf>=4.0"
    """

    extensions = (".pdf",)

    def extract(self, path: Path) -> str:
        try:
            from pypdf import PdfReader  # type: ignore[import]
        except ImportError as exc:
            raise ImportError(
                "pypdf is required for PDF ingestion.  "
                'Install it with:  pip install "pypdf>=4.0"'
            ) from exc

        reader = PdfReader(str(path))
        pages: list[str] = []
        for page in reader.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                pages.append(page_text)

        if not pages:
            raise ValueError(
                f"No extractable text found in '{path.name}'.  "
                "The PDF may be image-based (scanned) — consider running OCR first."
            )

        return "\n\n".join(pages)


# ── PowerPoint ────────────────────────────────────────────────────────────────

class PptxExtractor(ExtractorBase):
    """
    Extract text from PowerPoint (.pptx) files using ``python-pptx``.

    Install:  pip install "python-pptx>=1.0"

    Text is extracted in slide order: title → body → notes.
    """

    extensions = (".pptx", ".ppt")

    def extract(self, path: Path) -> str:
        if path.suffix.lower() == ".ppt":
            raise ValueError(
                "Legacy .ppt (PowerPoint 97-2003) is not supported.  "
                "Please save the file as .pptx and retry."
            )

        try:
            from pptx import Presentation  # type: ignore[import]
            from pptx.util import Pt  # noqa: F401 — confirms package is installed
        except ImportError as exc:
            raise ImportError(
                "python-pptx is required for PowerPoint ingestion.  "
                'Install it with:  pip install "python-pptx>=1.0"'
            ) from exc

        prs = Presentation(str(path))
        slides: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []

            # Slide title (if present)
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text_frame.text.strip()
                if title_text:
                    parts.append(f"[Slide {slide_num} Title] {title_text}")

            # All text frames (body, text boxes)
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue  # already captured above
                if shape.has_text_frame:
                    body_text = shape.text_frame.text.strip()
                    if body_text:
                        parts.append(body_text)

            # Speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"[Notes] {notes_text}")

            if parts:
                slides.append("\n".join(parts))

        if not slides:
            raise ValueError(
                f"No extractable text found in '{path.name}'.  "
                "The presentation may contain only images."
            )

        return "\n\n".join(slides)


# ── Registry & dispatcher ─────────────────────────────────────────────────────

_EXTRACTORS: list[ExtractorBase] = [
    TextExtractor(),
    PdfExtractor(),
    PptxExtractor(),
]

# Map extension → extractor (populated at import time)
_EXT_MAP: dict[str, ExtractorBase] = {}
for _ext in _EXTRACTORS:
    for _e in _ext.extensions:
        _EXT_MAP[_e] = _ext


def get_extractor(path: Path) -> ExtractorBase:
    """
    Return the appropriate extractor for *path* based on its file extension.

    Falls back to ``TextExtractor`` for unknown extensions (preserving the
    existing behaviour for plain-text files with unusual extensions).
    """
    ext = path.suffix.lower()
    return _EXT_MAP.get(ext, TextExtractor())


def extract_text(path: Path, enable_rich_formats: bool = True) -> str:
    """
    Extract plain text from *path*, routing to the correct extractor.

    Parameters
    ----------
    path               : Path to the document.
    enable_rich_formats: When ``False``, only ``.txt``-compatible files are
                         accepted.  Attempting to ingest a PDF or PPTX raises
                         ``ValueError`` — useful to hard-gate the old behaviour.

    Returns
    -------
    Plain-text string suitable for chunking/embedding.
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"Document not found: {path}")

    extractor = get_extractor(path)

    if not enable_rich_formats and not isinstance(extractor, TextExtractor):
        raise ValueError(
            f"Rich-format ingestion is disabled (enable_rich_formats=False).  "
            f"Cannot ingest '{path.suffix}' file '{path.name}'.  "
            "Set ingestion.enable_rich_formats: true in config.yaml to enable."
        )

    return extractor.extract(path)
